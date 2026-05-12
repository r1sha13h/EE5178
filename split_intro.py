"""Split Intro.pptx into two equal halves by deleting the unwanted slides
from a copy of the source. This preserves images and all relationships."""
import shutil
from pptx import Presentation

SRC = "Intro.pptx"
OUT1 = "Intro Part 1.pptx"
OUT2 = "Intro Part 2.pptx"

R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def delete_slides(pptx_path, keep_indices):
    """Open pptx_path and delete every slide whose original index is not in keep_indices."""
    prs = Presentation(pptx_path)
    sldIdLst = prs.slides._sldIdLst
    sldId_elements = list(sldIdLst)  # ordered list of <p:sldId> elements

    # iterate from the end so removal doesn't shift indices we still need
    for i in reversed(range(len(sldId_elements))):
        if i in keep_indices:
            continue
        sldId = sldId_elements[i]
        rId = sldId.get(R_NS)
        # drop the relationship and the underlying slide part
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    prs.save(pptx_path)


def main():
    prs = Presentation(SRC)
    n = len(prs.slides)
    print(f"Total slides: {n}")
    mid = n // 2
    part1 = set(range(0, mid))
    part2 = set(range(mid, n))
    print(f"Part 1: slides 1..{mid} ({len(part1)} slides)")
    print(f"Part 2: slides {mid+1}..{n} ({len(part2)} slides)")

    shutil.copyfile(SRC, OUT1)
    shutil.copyfile(SRC, OUT2)
    delete_slides(OUT1, part1)
    delete_slides(OUT2, part2)
    print(f"Wrote: {OUT1}")
    print(f"Wrote: {OUT2}")


if __name__ == "__main__":
    main()
